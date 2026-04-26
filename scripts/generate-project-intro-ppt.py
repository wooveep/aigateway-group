#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parent.parent
VERSION = "1.1.0"
OUT_DIR = REPO_ROOT / f"out/docs/{VERSION}"
ASSET_DIR = OUT_DIR / "assets/pptx-images"

IMG_HERO = ASSET_DIR / "ali-blue-hero.png"
IMG_ARCH = ASSET_DIR / "ali-blue-architecture.png"
IMG_SOLUTION = ASSET_DIR / "ali-blue-solution.png"
IMG_FLOW = ASSET_DIR / "ali-blue-flow.png"

W = Inches(13.333)
H = Inches(7.5)

BLUE = RGBColor(0, 86, 214)
DEEP_BLUE = RGBColor(0, 45, 156)
CYAN = RGBColor(0, 170, 255)
LIGHT_BLUE = RGBColor(232, 244, 255)
PALE_BLUE = RGBColor(245, 250, 255)
INK = RGBColor(18, 38, 63)
MUTED = RGBColor(91, 111, 136)
WHITE = RGBColor(255, 255, 255)


def add_text(slide, text, x, y, w, h, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bg(slide, color=PALE_BLUE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_image(slide, path: Path, x, y, w, h):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def add_title(slide, text, subtitle=""):
    add_text(slide, text, 0.62, 0.38, 7.8, 0.55, 28, True, DEEP_BLUE)
    if subtitle:
        add_text(slide, subtitle, 0.64, 0.98, 9.8, 0.38, 12, False, MUTED)


def card(slide, title, body, x, y, w, h, fill=WHITE, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(182, 216, 248)
    shape.line.width = Pt(1)
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.32, 13.5, True, accent)
    add_text(slide, body, x + 0.18, y + 0.56, w - 0.36, h - 0.68, 10.5, False, INK)
    return shape


def pill(slide, text, x, y, w, color=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text(slide, text, x + 0.08, y + 0.07, w - 0.16, 0.18, 9.5, True, WHITE, PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = CYAN
    line.line.width = Pt(1.6)
    line.line.end_arrowhead = True


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    add_image(slide, IMG_HERO, 0, 0, 13.333, 7.5)
    add_text(slide, "AIGateway", 0.72, 1.35, 4.8, 0.72, 42, True, DEEP_BLUE)
    add_text(slide, "企业级 AI Gateway 产品方案", 0.76, 2.14, 5.4, 0.42, 20, True, BLUE)
    add_text(slide, "统一接入大模型 API、MCP 工具与智能体应用，提供治理、计费、安全与观测闭环", 0.78, 2.72, 5.7, 0.62, 13, False, INK)
    for i, item in enumerate(["统一接入", "模型治理", "MCP 工具", "计费运营"]):
        pill(slide, item, 0.82 + i * 1.16, 5.88, 0.92, BLUE if i % 2 == 0 else CYAN)


def what_for(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    add_title(slide, "这个项目可以用来做什么", "把企业内分散的模型、工具、智能体和调用方统一纳入一套 AI 网关平台")
    cards = [
        ("统一模型入口", "对接 OpenAI、Anthropic、通义、Kimi、智谱等模型服务，向业务侧提供统一调用入口。"),
        ("统一工具入口", "把 MCP Server 和工具能力资产化，纳入授权、审计、运行时治理和门户展示。"),
        ("统一用户入口", "业务用户在 Portal 查看模型、申请 API Key、看账单、使用智能体和开放平台能力。"),
        ("统一治理入口", "管理员在 Console 管理组织、模型资产、AI Route、限流、脱敏、安全和系统任务。"),
        ("统一运营闭环", "请求级用量、钱包流水、计费、统计和观测数据形成平台运营闭环。"),
    ]
    for idx, (head, body) in enumerate(cards):
        card(slide, head, body, 0.78 + (idx % 3) * 4.1, 1.72 + (idx // 3) * 2.12, 3.55, 1.52)
    add_text(slide, "结果：业务只接一个 AI Gateway，平台集中完成模型接入、权限、额度、计费、安全与可观测。", 0.9, 6.35, 11.5, 0.36, 14, True, DEEP_BLUE)


def architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    add_image(slide, IMG_ARCH, 0.18, 1.08, 12.95, 5.95)
    add_title(slide, "总体架构", "参考 AI Gateway 的统一代理思路：应用与 Agent 只面对统一网关，后端连接模型、MCP 工具和治理体系")
    card(slide, "接入层", "业务应用 / Agent / Portal 用户 / Open Platform 调用方", 0.62, 1.65, 2.72, 1.0, fill=RGBColor(250, 253, 255))
    card(slide, "网关层", "Higress Gateway 执行鉴权、路由、协议转换、限流、脱敏和插件链。", 4.72, 2.25, 3.08, 1.18, fill=RGBColor(250, 253, 255))
    card(slide, "资源层", "模型 Provider、MCP Server、智能体目录、Wasm 插件和外部 API。", 9.35, 1.65, 2.95, 1.0, fill=RGBColor(250, 253, 255))
    card(slide, "控制与运营", "Console 配置资产和策略；Portal 承接用户自助、账务、用量和智能体入口。", 3.4, 5.58, 6.5, 0.95, fill=RGBColor(250, 253, 255))


def solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    add_title(slide, "产品方案", "Console 做平台治理，Portal 做用户自助，Gateway 做运行时执行")
    add_image(slide, IMG_SOLUTION, 0.25, 1.2, 12.8, 5.78)
    card(slide, "Console 控制面", "模型资产、Provider、AI Route、MCP、智能体目录、组织与等级、运行时策略。", 0.78, 5.55, 3.5, 1.05)
    card(slide, "AIGateway 运行时", "统一代理模型和工具请求，执行 key-auth、ai-quota、ai-token-ratelimit、ai-sensitive、ai-proxy。", 4.86, 5.55, 3.72, 1.05)
    card(slide, "Portal 用户门户", "模型广场、智能体广场、API Key、账单、发票、开放平台和 AI 对话。", 9.15, 5.55, 3.38, 1.05)


def flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    add_title(slide, "业务逻辑", "从平台配置到用户调用，再回流到计费与运营分析")
    add_image(slide, IMG_FLOW, 0.15, 1.2, 13.0, 5.6)
    steps = [
        ("1. 接入资产", "配置模型 Provider、模型绑定、MCP Server 和智能体目录。"),
        ("2. 配置策略", "设置授权对象、等级策略、AI Route、fallback、限流和脱敏。"),
        ("3. 用户自助", "用户查看模型/智能体，申请 API Key，选择可用能力。"),
        ("4. 网关执行", "统一鉴权、计量、路由、协议转换、工具调用和插件处理。"),
        ("5. 运营回流", "请求明细、账单、钱包流水、统计和观测数据回到平台。"),
    ]
    for idx, (head, body) in enumerate(steps):
        card(slide, head, body, 0.55 + idx * 2.52, 5.72, 2.18, 0.88, fill=RGBColor(250, 253, 255), accent=DEEP_BLUE)


def capabilities(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    add_title(slide, "核心功能", "围绕 AI Gateway 场景，把模型、工具、用户、策略和账务放到同一平台治理")
    items = [
        ("多模型统一接入", "统一管理 Provider、模型目录、模型绑定、协议和推荐入口。"),
        ("协议转换与路由", "OpenAI / Anthropic / Original 协议收口，AI Route 支持目标模型与 fallback。"),
        ("MCP 与智能体", "MCP Server 资产化，智能体目录连接工具、授权和 Portal 展示。"),
        ("API Key 与授权", "按 consumer、department、user_level 展开权限，支持高等级继承低等级。"),
        ("额度与限流", "支持金额额度、RPM/TPM 投影、Token 限流和 API Key 级运行时控制。"),
        ("内容安全", "AI 脱敏、敏感内容规则、插件链治理和安全审计。"),
        ("计费运营", "请求级明细、钱包流水、账单、发票和开放平台用量统计。"),
        ("可观测与扩展", "Prometheus 观测、系统任务、Wasm 插件分发和可扩展运行时能力。"),
    ]
    for idx, (head, body) in enumerate(items):
        x = 0.7 + (idx % 4) * 3.1
        y = 1.62 + (idx // 4) * 2.05
        card(slide, head, body, x, y, 2.72, 1.45)


def scenarios(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    add_title(slide, "典型使用场景", "适合企业把 AI 能力从“单点接入”升级为“平台化运营”")
    scenarios_data = [
        ("企业统一大模型出口", "多业务线通过统一网关调用不同供应商模型，平台集中做鉴权、路由、限流、审计和成本管理。"),
        ("内部智能体平台", "将 MCP 工具与智能体目录开放给员工或客户，按部门和等级控制可见性与调用范围。"),
        ("商业化开放平台", "对外提供 API Key、模型广场、账单、发票、用量统计和调用明细，支持按组织运营。"),
        ("多模型治理与降本", "通过模型资产、价格版本、路由策略、fallback 和用量分析，平衡成本、稳定性和体验。"),
    ]
    for idx, (head, body) in enumerate(scenarios_data):
        x = 0.85 + (idx % 2) * 5.92
        y = 1.65 + (idx // 2) * 2.28
        card(slide, head, body, x, y, 5.25, 1.65, fill=RGBColor(250, 253, 255), accent=DEEP_BLUE)
    add_text(slide, "一句话：AIGateway 让企业可以像管理云资源一样管理 AI 模型、工具、智能体、调用方和成本。", 0.92, 6.52, 11.6, 0.36, 15, True, BLUE)


def build(output: Path):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    cover(prs)
    what_for(prs)
    architecture(prs)
    solution(prs)
    flow(prs)
    capabilities(prs)
    scenarios(prs)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate AIGateway product introduction PPT.")
    parser.add_argument(
        "--output",
        default=str(OUT_DIR / "aigateway-project-introduction-1.1.0.pptx"),
        help="Output PPTX path.",
    )
    args = parser.parse_args()
    build(Path(args.output).resolve())
    print(args.output)


if __name__ == "__main__":
    main()
